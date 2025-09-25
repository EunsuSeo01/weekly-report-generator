"""
시스템 아키텍처 PPT 생성 스크립트
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_architecture_ppt():
    # 새 프레젠테이션 생성
    prs = Presentation()

    # 슬라이드 1: 백엔드 + 프론트엔드 시스템 아키텍처
    slide1_layout = prs.slide_layouts[5]  # 빈 슬라이드
    slide1 = prs.slides.add_slide(slide1_layout)

    # 제목 추가
    title_shape = slide1.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    title_frame = title_shape.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "AI 기반 주간업무 보고서 생성 시스템 아키텍처"
    title_p.font.size = Pt(28)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(31, 119, 180)
    title_p.alignment = PP_ALIGN.CENTER

    # 프론트엔드 박스
    frontend_shape = slide1.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(2.5), Inches(2))
    frontend_frame = frontend_shape.text_frame
    frontend_frame.margin_left = Inches(0.1)
    frontend_frame.margin_right = Inches(0.1)
    frontend_frame.margin_top = Inches(0.1)
    frontend_frame.margin_bottom = Inches(0.1)

    p1 = frontend_frame.paragraphs[0]
    p1.text = "Frontend Layer"
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.alignment = PP_ALIGN.CENTER

    p2 = frontend_frame.add_paragraph()
    p2.text = "\n• Vue.js + Vite"
    p2.font.size = Pt(12)

    p3 = frontend_frame.add_paragraph()
    p3.text = "• Streamlit Dashboard"
    p3.font.size = Pt(12)

    p4 = frontend_frame.add_paragraph()
    p4.text = "• 사용자 인터페이스"
    p4.font.size = Pt(12)

    # 백엔드 박스
    backend_shape = slide1.shapes.add_textbox(Inches(4), Inches(1.5), Inches(2.5), Inches(2))
    backend_frame = backend_shape.text_frame
    backend_frame.margin_left = Inches(0.1)
    backend_frame.margin_right = Inches(0.1)
    backend_frame.margin_top = Inches(0.1)
    backend_frame.margin_bottom = Inches(0.1)

    b1 = backend_frame.paragraphs[0]
    b1.text = "Backend API Server"
    b1.font.size = Pt(16)
    b1.font.bold = True
    b1.alignment = PP_ALIGN.CENTER

    b2 = backend_frame.add_paragraph()
    b2.text = "\n• FastAPI (포트 8001)"
    b2.font.size = Pt(12)

    b3 = backend_frame.add_paragraph()
    b3.text = "• psycopg2 DB 연결"
    b3.font.size = Pt(12)

    b4 = backend_frame.add_paragraph()
    b4.text = "• CORS 미들웨어"
    b4.font.size = Pt(12)

    # 데이터베이스 박스
    db_shape = slide1.shapes.add_textbox(Inches(7.5), Inches(1.5), Inches(2), Inches(2))
    db_frame = db_shape.text_frame
    db_frame.margin_left = Inches(0.1)
    db_frame.margin_right = Inches(0.1)
    db_frame.margin_top = Inches(0.1)
    db_frame.margin_bottom = Inches(0.1)

    d1 = db_frame.paragraphs[0]
    d1.text = "Database"
    d1.font.size = Pt(16)
    d1.font.bold = True
    d1.alignment = PP_ALIGN.CENTER

    d2 = db_frame.add_paragraph()
    d2.text = "\n• PostgreSQL"
    d2.font.size = Pt(12)

    d3 = db_frame.add_paragraph()
    d3.text = "• pgvector 확장"
    d3.font.size = Pt(12)

    d4 = db_frame.add_paragraph()
    d4.text = "• 벡터 검색"
    d4.font.size = Pt(12)

    # AI 서비스 박스
    ai_shape = slide1.shapes.add_textbox(Inches(4), Inches(4.2), Inches(2.5), Inches(1.8))
    ai_frame = ai_shape.text_frame
    ai_frame.margin_left = Inches(0.1)
    ai_frame.margin_right = Inches(0.1)
    ai_frame.margin_top = Inches(0.1)
    ai_frame.margin_bottom = Inches(0.1)

    a1 = ai_frame.paragraphs[0]
    a1.text = "AI Services"
    a1.font.size = Pt(16)
    a1.font.bold = True
    a1.alignment = PP_ALIGN.CENTER

    a2 = ai_frame.add_paragraph()
    a2.text = "\n• LangChain"
    a2.font.size = Pt(12)

    a3 = ai_frame.add_paragraph()
    a3.text = "• OpenAI GPT-4o"
    a3.font.size = Pt(12)

    a4 = ai_frame.add_paragraph()
    a4.text = "• 보고서 자동 생성"
    a4.font.size = Pt(12)

    # 데이터 소스 박스들
    sources = [
        ("Slack", "메시지/채널"),
        ("Notion", "문서/페이지"),
        ("OneDrive", "파일 활동"),
        ("Outlook", "이메일/캘린더")
    ]

    for i, (source, desc) in enumerate(sources):
        x_pos = 0.5 + i * 2.2
        source_shape = slide1.shapes.add_textbox(Inches(x_pos), Inches(6.2), Inches(2), Inches(1.2))
        source_frame = source_shape.text_frame
        source_frame.margin_left = Inches(0.1)
        source_frame.margin_right = Inches(0.1)
        source_frame.margin_top = Inches(0.1)
        source_frame.margin_bottom = Inches(0.1)

        s1 = source_frame.paragraphs[0]
        s1.text = source
        s1.font.size = Pt(14)
        s1.font.bold = True
        s1.alignment = PP_ALIGN.CENTER

        s2 = source_frame.add_paragraph()
        s2.text = f"\n{desc}"
        s2.font.size = Pt(10)
        s2.alignment = PP_ALIGN.CENTER

    # 슬라이드 2: API 서비스 플로우 흐름도
    slide2_layout = prs.slide_layouts[5]  # 빈 슬라이드
    slide2 = prs.slides.add_slide(slide2_layout)

    # 제목 추가
    title2_shape = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    title2_frame = title2_shape.text_frame
    title2_p = title2_frame.paragraphs[0]
    title2_p.text = "API 서비스 플로우 흐름도"
    title2_p.font.size = Pt(28)
    title2_p.font.bold = True
    title2_p.font.color.rgb = RGBColor(31, 119, 180)
    title2_p.alignment = PP_ALIGN.CENTER

    # 플로우 단계들
    flow_steps = [
        ("1. 사용자 요청", "user_id, start_date, end_date"),
        ("2. 데이터 수집", "4개 플랫폼 병렬 조회"),
        ("3. 타임라인 생성", "시간순 정렬 및 통합"),
        ("4. AI 보고서 생성", "GPT-4o + 프롬프트 템플릿"),
        ("5. DB 저장", "report 테이블에 결과 저장"),
        ("6. 응답 반환", "JSON 형태 결과 반환")
    ]

    # 세로 플로우로 배치
    for i, (step, desc) in enumerate(flow_steps):
        y_pos = 1.5 + i * 0.9

        # 단계 박스
        step_shape = slide2.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(3), Inches(0.7))
        step_frame = step_shape.text_frame
        step_frame.margin_left = Inches(0.1)
        step_frame.margin_right = Inches(0.1)
        step_frame.margin_top = Inches(0.1)
        step_frame.margin_bottom = Inches(0.1)

        s1 = step_frame.paragraphs[0]
        s1.text = step
        s1.font.size = Pt(14)
        s1.font.bold = True
        s1.font.color.rgb = RGBColor(31, 119, 180)

        # 설명 박스
        desc_shape = slide2.shapes.add_textbox(Inches(4.5), Inches(y_pos), Inches(4.5), Inches(0.7))
        desc_frame = desc_shape.text_frame
        desc_frame.margin_left = Inches(0.1)
        desc_frame.margin_right = Inches(0.1)
        desc_frame.margin_top = Inches(0.1)
        desc_frame.margin_bottom = Inches(0.1)

        d1 = desc_frame.paragraphs[0]
        d1.text = desc
        d1.font.size = Pt(12)

        # 화살표 (마지막 단계 제외)
        if i < len(flow_steps) - 1:
            arrow_shape = slide2.shapes.add_textbox(Inches(2.3), Inches(y_pos + 0.7), Inches(0.4), Inches(0.2))
            arrow_frame = arrow_shape.text_frame
            arrow_p = arrow_frame.paragraphs[0]
            arrow_p.text = "↓"
            arrow_p.font.size = Pt(16)
            arrow_p.font.bold = True
            arrow_p.font.color.rgb = RGBColor(31, 119, 180)
            arrow_p.alignment = PP_ALIGN.CENTER

    # API 엔드포인트 정보
    api_shape = slide2.shapes.add_textbox(Inches(1), Inches(7.2), Inches(8), Inches(1))
    api_frame = api_shape.text_frame
    api_frame.margin_left = Inches(0.1)
    api_frame.margin_right = Inches(0.1)
    api_frame.margin_top = Inches(0.1)
    api_frame.margin_bottom = Inches(0.1)

    api1 = api_frame.paragraphs[0]
    api1.text = "핵심 API 엔드포인트"
    api1.font.size = Pt(14)
    api1.font.bold = True

    api2 = api_frame.add_paragraph()
    api2.text = "• GET /api/user-timeline/{user_id} - 사용자별 통합 타임라인 조회"
    api2.font.size = Pt(11)

    api3 = api_frame.add_paragraph()
    api3.text = "• POST /reports/weekly - AI 기반 주간 보고서 생성 및 저장"
    api3.font.size = Pt(11)

    api4 = api_frame.add_paragraph()
    api4.text = "• GET /health - 서버 상태 확인"
    api4.font.size = Pt(11)

    # PPT 파일 저장
    prs.save('system_architecture.pptx')
    return 'system_architecture.pptx'

if __name__ == "__main__":
    filename = create_architecture_ppt()
    print(f"PPT 파일이 생성되었습니다: {filename}")